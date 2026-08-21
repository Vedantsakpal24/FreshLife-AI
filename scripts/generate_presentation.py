"""
FreshLife AI - 12 Slide Professional Presentation Generator
Generates FreshLife_AI_Presentation.pptx using python-pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Color Palette
    DARK_BG = RGBColor(15, 30, 20)        # Deep Forest Dark
    LIGHT_BG = RGBColor(248, 250, 248)    # Crisp Off-White
    PRIMARY_GREEN = RGBColor(34, 197, 94) # Vibrant Emerald
    DARK_GREEN = RGBColor(20, 83, 45)     # Dark Forest Text
    CARD_BG = RGBColor(255, 255, 255)     # Clean Card
    CARD_BORDER = RGBColor(220, 235, 225) # Subtle Border
    TEXT_MAIN = RGBColor(30, 41, 59)      # Slate 800
    TEXT_MUTED = RGBColor(100, 116, 139)  # Slate 500
    ACCENT_AMBER = RGBColor(245, 158, 11) # Warning/Alert
    WHITE = RGBColor(255, 255, 255)

    def add_header(slide, title_text, category_text="FRESHLIFE AI"):
        # Category pill
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(10), Inches(0.35))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = PRIMARY_GREEN
        
        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.5), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = DARK_GREEN

    def add_card(slide, left, top, width, height, title, items, badge=""):
        # Background card shape
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = CARD_BORDER
        shape.line.width = Pt(1.2)
        
        # Card content
        pad = Inches(0.25)
        tb = slide.shapes.add_textbox(left + pad, top + pad, width - (pad * 2), height - (pad * 2))
        tf = tb.text_frame
        tf.word_wrap = True
        
        if badge:
            p_badge = tf.paragraphs[0]
            p_badge.text = badge.upper()
            p_badge.font.size = Pt(9)
            p_badge.font.bold = True
            p_badge.font.color.rgb = PRIMARY_GREEN
            p_title = tf.add_paragraph()
        else:
            p_title = tf.paragraphs[0]
            
        p_title.text = title
        p_title.font.size = Pt(16)
        p_title.font.bold = True
        p_title.font.color.rgb = DARK_GREEN
        p_title.space_after = Pt(10)
        
        for item in items:
            p = tf.add_paragraph()
            if isinstance(item, tuple):
                p.text = f"• {item[0]}: "
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = TEXT_MAIN
                run = p.add_run()
                run.text = item[1]
                run.font.bold = False
                run.font.color.rgb = TEXT_MUTED
            else:
                p.text = f"• {item}"
                p.font.size = Pt(12)
                p.font.color.rgb = TEXT_MAIN
            p.space_after = Pt(6)

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (Dark Theme)
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_slide_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = DARK_BG
    bg1.line.fill.background()

    # Title content
    tb1 = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(4.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "🌿 FRESHLIFE AI"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_GREEN
    p.space_after = Pt(12)
    
    p = tf1.add_paragraph()
    p.text = "Smart Fruit & Vegetable Freshness Detection & Shelf-Life Estimation"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(16)
    
    p = tf1.add_paragraph()
    p.text = "An end-to-end computer vision and environmental intelligence platform designed to combat food waste through accurate freshness scoring and adaptive shelf-life predictions."
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(200, 220, 205)
    p.space_after = Pt(28)
    
    p = tf1.add_paragraph()
    p.text = "Full Stack Architecture  |  Ultralytics YOLO  |  FastAPI  |  Next.js 16  |  MongoDB Atlas"
    p.font.size = Pt(12)
    p.font.color.rgb = PRIMARY_GREEN

    # -------------------------------------------------------------
    # SLIDE 2: Problem Statement & Motivation
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_slide_layout)
    add_header(s2, "The Problem: Perishable Food Spoilage & Waste", "Problem & Motivation")
    add_card(s2, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8), 
             "Global Food Waste Crisis", [
                 ("1.3 Billion Tons", "Global annual food waste, with fruits and vegetables having the highest wastage rate (~45%)."),
                 ("Economic Loss", "Billions lost annually across households, grocery stores, and supply chain operators."),
                 ("Environmental Impact", "Decomposing food in landfills contributes heavily to methane emissions.")
             ], "CRITICAL ISSUE")
    
    add_card(s2, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8), 
             "Consumer & Retail Challenges", [
                 ("Subjective Visual Inspection", "Consumers struggle to judge internal fruit freshness accurately."),
                 ("Static Expiration Dates", "Produce does not come with fixed expiry dates; it degrades dynamically."),
                 ("Suboptimal Storage", "Lack of knowledge on temperature and storage condition impacts.")
             ], "DAILY FRICTION")
             
    add_card(s2, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.8), 
             "The Solution: FreshLife AI", [
                 ("Visual AI Scanning", "Instant classification of produce and ripeness state from a single photo."),
                 ("Environmental Adaptation", "Considers ambient temperature, humidity, and storage type."),
                 ("Actionable Guidance", "Provides estimated shelf-life days and tailored consumption advice.")
             ], "OUR APPROACH")

    # -------------------------------------------------------------
    # SLIDE 3: Project Vision & Core Objectives
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_slide_layout)
    add_header(s3, "Project Objectives & Key Goals", "Project Scope")
    add_card(s3, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.3),
             "1. High-Accuracy Visual Classification", [
                 "Train and deploy a lightweight YOLO model (`best.pt`) supporting 24 distinct produce conditions.",
                 "Achieve sub-second inference on CPU for fast real-time analysis."
             ], "OBJECTIVE")

    add_card(s3, Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.3),
             "2. Dynamic Environmental Shelf-Life Engine", [
                 "Combine visual condition with storage parameters (temp, humidity, storage method).",
                 "Deliver realistic Min–Max shelf life day ranges rather than static guesses."
             ], "OBJECTIVE")

    add_card(s3, Inches(0.8), Inches(4.4), Inches(5.6), Inches(2.3),
             "3. Explainable AI & Consumer Actionability", [
                 "Provide transparent reasoning behind every score to build user trust.",
                 "Generate proactive storage recommendations (e.g., refrigerate, consume immediately)."
             ], "OBJECTIVE")

    add_card(s3, Inches(6.8), Inches(4.4), Inches(5.7), Inches(2.3),
             "4. Cross-Platform Accessibility", [
                 "Develop responsive Next.js frontend with live camera streaming and mobile support.",
                 "Cloud persistence with MongoDB Atlas for scan history tracking."
             ], "OBJECTIVE")

    # -------------------------------------------------------------
    # SLIDE 4: System Architecture & Tech Stack
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_slide_layout)
    add_header(s4, "Full-Stack System Architecture", "Technical Stack")
    add_card(s4, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8),
             "Frontend Layer", [
                 ("Framework", "Next.js 16 (App Router) & React 19"),
                 ("Styling", "Tailwind CSS v4 (Modern Design)"),
                 ("Language", "TypeScript (Strict Types)"),
                 ("Capabilities", "Camera streaming, file upload, reactive progress bar, score gauges")
             ], "CLIENT SIDE")

    add_card(s4, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8),
             "ML & Backend Layer", [
                 ("Engine", "FastAPI (Asynchronous Python)"),
                 ("Model Runtime", "Ultralytics YOLO (`best.pt`)"),
                 ("Image Processing", "Pillow (PIL), NumPy"),
                 ("Lifecycle", "Single startup model loading with persistent memory caching")
             ], "CORE BACKEND")

    add_card(s4, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.8),
             "Database & Cloud", [
                 ("Database", "MongoDB Atlas Cloud Cluster"),
                 ("ODM", "Mongoose with cached connection pools"),
                 ("Deployment", "Vercel (Frontend) & Cloud Container (Backend)"),
                 ("Security", "HTTPS encryption for camera permissions & TLS connections")
             ], "STORAGE & CLOUD")

    # -------------------------------------------------------------
    # SLIDE 5: Computer Vision & AI Model (best.pt)
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_slide_layout)
    add_header(s5, "Computer Vision Model: YOLO Architecture", "Machine Learning")
    add_card(s5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8),
             "Model Specifications & Classes", [
                 ("Architecture", "Custom YOLO Deep Neural Network (`best.pt`)"),
                 ("Model Footprint", "Lightweight ~3.2 MB (Ideal for fast CPU edge inference)"),
                 ("Supported Fruits", "Banana, Orange, Papaya, Pineapple, Apple"),
                 ("Supported Vegetables", "Bittermelon, Cucumber, Eggplant, Tomato"),
                 ("Condition Classes", "Fresh, Semi-Fresh, Rotten (Total 24 combinations)"),
                 ("Inference Latency", "< 250ms per scan on standard CPU")
             ], "MODEL SPECS")

    add_card(s5, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8),
             "Inference Pipeline Flow", [
                 ("1. Input Ingestion", "Image decoded in-memory as RGB via PIL without temp disk writes."),
                 ("2. Forward Pass", "Model extracts top probability class and confidence score."),
                 ("3. Token Splitting", "Decodes class string (e.g. `fresh_banana` → Produce: Banana, Condition: Fresh)."),
                 ("4. Freshness Scoring", "Computes calibrated Freshness Score (0–100 scale)."),
                 ("5. Dual Compatibility", "Auto-handles both YOLO classification (`probs`) and detection (`boxes`).")
             ], "PIPELINE LOGIC")

    # -------------------------------------------------------------
    # SLIDE 6: Dynamic Environmental Shelf-Life Engine
    # -------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_slide_layout)
    add_header(s6, "Environmental Shelf-Life Engine", "Algorithmic Logic")
    add_card(s6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8),
             "Multi-Factor Environmental Model", [
                 ("Base Estimation", "Establishes baseline lifespan for detected produce condition at 25°C."),
                 ("Temperature Impact", "Applies exponential factor for deviations above/below baseline temperature: per-degree degradation rate."),
                 ("Humidity Adjustments", "Penalizes extreme moisture (>80% risk of mold) or dryness (<30% risk of desiccation)."),
                 ("Storage Multipliers", "Room Temperature (1.0x), Pantry (1.2x), Refrigerator (2.5x), Freezer (6.0x).")
             ], "CALCULATION FORMULA")

    add_card(s6, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8),
             "Real-World Prediction Example", [
                 ("Scenario 1: Fresh Banana @ Room Temp (25°C)", "Baseline: 5–7 Days. Result: 5–7 Days."),
                 ("Scenario 2: Fresh Banana in Refrigerator (4°C)", "Temp + Storage Adjustment. Result: 12–15 Days."),
                 ("Scenario 3: Semi-Fresh Tomato @ High Heat (34°C)", "Accelerated spoilage factor applied. Result: 1–2 Days."),
                 ("Scenario 4: Rotten Fruit Detected", "Shelf life automatically clamped to 0 Days with safety alert.")
             ], "CASE STUDIES")

    # -------------------------------------------------------------
    # SLIDE 7: User Experience & Scan Workflow
    # -------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_slide_layout)
    add_header(s7, "User Experience & Scanning Workflow", "UX & Interface")
    add_card(s7, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8),
             "1. Input Capture", [
                 ("Live Viewfinder", "In-browser camera stream with visual fruit framing reticle."),
                 ("Camera Switcher", "Seamless toggle between front and rear cameras on mobile."),
                 ("File & Native Camera", "One-click access to native phone camera and file uploads.")
             ], "STEP 1")

    add_card(s7, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8),
             "2. Environmental Inputs", [
                 ("Temperature Slider", "User inputs ambient temperature (°C) with presets."),
                 ("Humidity Control", "Selects relative humidity percentage (%)."),
                 ("Storage Selector", "Chooses storage type (Room Temp, Fridge, Pantry, Freezer).")
             ], "STEP 2")

    add_card(s7, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.8),
             "3. Real-Time Feedback", [
                 ("Progressive Feedback", "Animated indicator: Uploading → Identifying → Estimating."),
                 ("Fast Response", "Instant redirect to detailed results page once scan completes."),
                 ("Error Recovery", "Clear guidance if image is blurry or unsupported.")
             ], "STEP 3")

    # -------------------------------------------------------------
    # SLIDE 8: Explainable AI & Results Dashboard
    # -------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_slide_layout)
    add_header(s8, "Explainable AI & Results Dashboard", "Results & Insights")
    add_card(s8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8),
             "Core Dashboard Components", [
                 ("Produce Identity Card", "Displays item name, category badge, and AI confidence level (e.g. 98%)."),
                 ("Freshness Meter", "Color-coded gauge from 0 to 100 (Green = Fresh, Amber = Semi-Fresh, Red = Rotten)."),
                 ("Shelf-Life Range", "Clear Min–Max days count with storage condition indicator."),
                 ("Environmental Summary", "Recap of temperature, humidity, and storage context used for the prediction.")
             ], "UI COMPONENTS")

    add_card(s8, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8),
             "Explainability & Recommendations", [
                 ("Transparent Reasoning", "Bullet list detailing every contributing factor (AI visual analysis + temperature impact + storage method)."),
                 ("Actionable Recommendations", "Actionable advice based on condition (e.g. 'Good condition. Consider refrigerating to extend freshness')."),
                 ("Food Safety Warnings", "Clear safety disclaimer reminding consumers to physically check produce before consumption.")
             ], "EXPLAINABILITY")

    # -------------------------------------------------------------
    # SLIDE 9: Database & History Management
    # -------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_slide_layout)
    add_header(s9, "Data Persistence & Scan History", "Database Architecture")
    add_card(s9, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8),
             "MongoDB Atlas Schema (`ProduceScan`)", [
                 ("Image Metadata", "Stores base64 / image URI along with scan timestamp."),
                 ("AI Predictions", "Produce type, condition, freshness score, confidence rating."),
                 ("Environmental Parameters", "Recorded temperature (°C), humidity (%), and storage mode."),
                 ("Lifespan Estimates", "Calculated minimum and maximum shelf-life days."),
                 ("Explanations & Advice", "Full list of explainability notes and recommendation text.")
             ], "SCHEMA DESIGN")

    add_card(s9, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8),
             "Scan History Feature", [
                 ("Historical Records", "Dedicated `/history` page displaying all past user analyses."),
                 ("Timeline View", "Chronological list with produce emoji, freshness badge, and date."),
                 ("Instant Recall", "Click any previous scan to revisit the complete interactive result dashboard."),
                 ("Connection Pooling", "Global Mongoose connection caching to prevent connection exhaustion in serverless environments.")
             ], "USER FEATURES")

    # -------------------------------------------------------------
    # SLIDE 10: Performance & Key Features Summary
    # -------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_slide_layout)
    add_header(s10, "Performance Benchmarks & Key Features", "Evaluation")
    add_card(s10, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8),
             "Speed & Efficiency", [
                 ("< 250ms ML Inference", "Fast image processing on CPU."),
                 ("Single Model Load", "Persistent model in memory avoids cold start penalties."),
                 ("Sub-Second Full Flow", "End-to-end API response under 800ms.")
             ], "PERFORMANCE")

    add_card(s10, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8),
             "Accuracy & Coverage", [
                 ("24 Classes", "Comprehensive coverage of 8 major produce varieties."),
                 ("3-Stage Ripeness", "Distinguishes fresh, semi-fresh, and spoiled states."),
                 ("Confidence Gate", "Rejects low-confidence/unsupported images gracefully.")
             ], "MODEL COVERAGE")

    add_card(s10, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.8),
             "Architecture Strengths", [
                 ("Decoupled Stack", "Independent Next.js frontend and FastAPI ML service."),
                 ("Vercel Ready", "Production-optimized build with 0 TypeScript/ESLint errors."),
                 ("Secure Camera Access", "HTTPS compliance for seamless browser camera permissions.")
             ], "ROBUSTNESS")

    # -------------------------------------------------------------
    # SLIDE 11: Deployment & Cloud Infrastructure
    # -------------------------------------------------------------
    s11 = prs.slides.add_slide(blank_slide_layout)
    add_header(s11, "Cloud Deployment & Production Setup", "Deployment")
    add_card(s11, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8),
             "Frontend on Vercel", [
                 ("Global Edge CDN", "High-performance distribution with sub-100ms global delivery."),
                 ("Automatic HTTPS", "Mandatory TLS certificates ensure camera API permissions work everywhere."),
                 ("Environment Variables", "Secure injection of `MONGODB_URI` and `ML_SERVER_URL`."),
                 ("Optimized Bundles", "Turbopack static pre-rendering with React Suspense boundaries.")
             ], "VERCEL HOSTING")

    add_card(s11, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8),
             "Backend & Database Cloud", [
                 ("FastAPI Container", "Deployable on Render, Railway, or Hugging Face Spaces."),
                 ("MongoDB Atlas", "Serverless cloud database with automated replication and backups."),
                 ("CORS & Proxying", "Next.js API route proxies multipart form data securely to ML server."),
                 ("Zero Local Storage Dependencies", "Entire pipeline processes files in-memory and in cloud DB.")
             ], "BACKEND & DATA")

    # -------------------------------------------------------------
    # SLIDE 12: Future Roadmap & Conclusion
    # -------------------------------------------------------------
    s12 = prs.slides.add_slide(blank_slide_layout)
    add_header(s12, "Future Roadmap & Conclusion", "Vision")
    add_card(s12, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8),
             "Future Enhancements", [
                 ("Multi-Item Object Detection", "Upgrade to multi-fruit bounding box detection for scanning whole fruit bowls."),
                 ("Smart Fridge IoT Integration", "Connect with camera-equipped smart refrigerators for automatic inventory monitoring."),
                 ("Recipe Recommendation Engine", "Suggest recipes to use up semi-fresh produce before it spoils."),
                 ("Mobile App Release", "Publish dedicated iOS/Android apps with offline caching.")
             ], "ROADMAP")

    add_card(s12, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8),
             "Project Summary", [
                 ("Impact", "Delivers a practical, accessible AI tool to reduce global household and retail food waste."),
                 ("Technical Milestone", "Seamlessly integrates custom YOLO computer vision with environmental shelf-life logic."),
                 ("Production Ready", "Complete full-stack system with modern UI, explainability, and cloud persistence."),
                 ("Thank You!", "Questions & Discussion.")
             ], "CONCLUSION")

    # Save presentation
    output_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "FreshLife_AI_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    create_deck()
